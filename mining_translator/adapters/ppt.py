"""PowerPoint adapter for .pptx files using python-pptx."""

import os
from .base import BaseAdapter, TextBlock


class PPTAdapter(BaseAdapter):
    supported_extensions = [".pptx"]

    def extract_texts(self, filepath: str) -> list[TextBlock]:
        from pptx import Presentation

        prs = Presentation(filepath)
        blocks = []

        for si, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if shape.has_text_frame:
                    full_text = shape.text_frame.text.strip()
                    if full_text:
                        blocks.append(TextBlock(
                            id=f"s{si}|{shape.shape_id}",
                            text=full_text,
                            context=f"Slide {si+1}",
                            meta={"slide": si, "shape_id": shape.shape_id},
                        ))

                if shape.has_table:
                    for ri, row in enumerate(shape.table.rows):
                        for ci, cell in enumerate(row.cells):
                            if cell.text.strip():
                                blocks.append(TextBlock(
                                    id=f"s{si}|{shape.shape_id}|t_r{ri}_c{ci}",
                                    text=cell.text,
                                    context=f"Slide {si+1}, table",
                                    meta={"slide": si, "shape_id": shape.shape_id, "row": ri, "col": ci},
                                ))

            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    blocks.append(TextBlock(
                        id=f"s{si}|notes",
                        text=notes,
                        context=f"Slide {si+1} notes",
                        meta={"slide": si, "type": "notes"},
                    ))

        return blocks

    def write_translated(self, filepath: str, blocks: list[TextBlock], output_path: str):
        from pptx import Presentation

        prs = Presentation(filepath)
        lookup = {b.id: b.translated for b in blocks}

        for si, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                shape_id = f"s{si}|{shape.shape_id}"
                if shape.has_text_frame and shape_id in lookup:
                    _replace_text_frame(shape.text_frame, lookup[shape_id])

                if shape.has_table:
                    for ri, row in enumerate(shape.table.rows):
                        for ci, cell in enumerate(row.cells):
                            cell_id = f"s{si}|{shape.shape_id}|t_r{ri}_c{ci}"
                            if cell_id in lookup:
                                cell.text = lookup[cell_id]

            notes_id = f"s{si}|notes"
            if slide.has_notes_slide and notes_id in lookup:
                notes_frame = slide.notes_slide.notes_text_frame
                _replace_text_frame(notes_frame, lookup[notes_id])

        prs.save(output_path)


def _replace_text_frame(tf, new_text: str):
    """Replace all text in a text frame."""
    paragraphs = new_text.split("\n")
    # Update existing paragraphs
    for i, para in enumerate(tf.paragraphs):
        if i < len(paragraphs):
            if para.runs:
                para.runs[0].text = paragraphs[i]
                for run in para.runs[1:]:
                    run.text = ""
            else:
                para.text = paragraphs[i]
        else:
            para.text = ""
    # Add new paragraphs if needed
    for extra in paragraphs[len(tf.paragraphs):]:
        p = tf.add_paragraph()
        p.text = extra
